from django_extensions.management.jobs import DailyJob

import requests
import plotly.graph_objects as go

import pandas as pd
import matplotlib.pyplot as plt
import numpy as np
from datetime import datetime, timedelta
import mysql.connector
from pptx import Presentation
from pptx.util import Inches
import xlsxwriter
import smtplib
import os

from paramiko import SSHClient, AutoAddPolicy
from scp import SCPClient

from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.mime.application import MIMEApplication

from apps.app.models import HEALTHCHECK

class Job(DailyJob):
    help = "Script d'automatisation du healthcheck journalier avec envoie par email !!!"

    def execute(self):
        
        mydb = mysql.connector.connect(
            host="10.39.231.23",
            user="ABOKOR",
            password="Abk_2021",
            database="TOPUPDJIB"
        )

        mycursor = mydb.cursor()

        """ La date et l'heure """
        time_in = "09:00:00"
        time_su = "09:00:00"
        datetime_today = datetime.now()
        date_today = datetime_today.strftime("%Y-%m-%d")
        timestamp_today = str(datetime.timestamp(datetime_today))
        datetime_yesterday = datetime.now() - timedelta(1)
        timestamp_yesterday = str(datetime.timestamp(datetime_yesterday))
        date_yesterday = datetime_yesterday.strftime("%Y-%m-%d")
        datetime_in = date_yesterday+" "+time_in
        datetime_su = date_today+" "+time_su

        """ Les requêtes """
        mycursor.execute("SELECT CASE WHEN qh.action_type IN (1, 2, 3, 4, 5) THEN 'USSD' WHEN qh.action_type IN (44, 55) THEN 'CRM' ELSE CONCAT('ACTION_TYPE = ',qh.action_type) END AS SOURCE, CASE WHEN qh.action_type = 1 THEN 'CHANGEMENT CODE PIN' WHEN qh.action_type = 2 THEN 'AJOUTER NUMERO AU REPERTOIRE' WHEN qh.action_type = 3 THEN 'SUPPRIMER NUMERO DU REPERTOIRE' WHEN qh.action_type = 4 THEN 'ACHAT FORFAIT' WHEN qh.action_type = 5 THEN 'RECHARGE BALANCE PRINCIPALE' WHEN qh.action_type = 44 THEN 'ACHAT FORFAIT' WHEN qh.action_type = 55 THEN 'RECHARGE BALANCE PRINCIPALE' ELSE CONCAT('ACTION_TYPE = ',qh.action_type) END AS OPERATION, SUBSTRING(qh.command_content,LOCATE('#',qh.command_content)+1) AS CODE_PACKAGE_OU_MONTANT_RECHARGE, COUNT(*) AS NOMBRE, pd.frais, sum(pd.frais) as Total_Per_Package FROM queue_hist qh, package_data pd WHERE qh.created_date > '"+datetime_in+"' AND qh.created_date < '"+datetime_su+"' AND qh.action_type IN (4, 5, 44, 55) and pd.code = SUBSTRING(qh.command_content,LOCATE('#',qh.command_content)+1) GROUP BY SOURCE, OPERATION, CODE_PACKAGE_OU_MONTANT_RECHARGE ORDER BY SOURCE, NOMBRE DESC;")
        request_success = mycursor.fetchall()

        mycursor.execute("SELECT CASE WHEN q.action_type IN (1, 2, 3, 4, 5) THEN 'USSD' WHEN q.action_type IN (44, 55) THEN 'CRM' ELSE CONCAT('ACTION_TYPE = ',q.action_type) END AS SOURCE, CASE WHEN q.action_type = 1 THEN 'CHANGEMENT CODE PIN' WHEN q.action_type = 2 THEN 'AJOUTER NUMERO AU REPERTOIRE' WHEN q.action_type = 3 THEN 'SUPPRIMER NUMERO DU REPERTOIRE' WHEN q.action_type = 4 THEN 'ACHAT FORFAIT' WHEN q.action_type = 5 THEN 'RECHARGE BALANCE PRINCIPALE' WHEN q.action_type = 44 THEN 'ACHAT FORFAIT' WHEN q.action_type = 55 THEN 'RECHARGE BALANCE PRINCIPALE' ELSE CONCAT('ACTION_TYPE = ',q.action_type) END AS OPERATION, q.error_description AS MESSAGE_ERREUR, COUNT(*) AS NOMBRE FROM queue q WHERE q.created_date > '"+datetime_in+"' AND q.created_date < '"+datetime_su+"' GROUP BY SOURCE, OPERATION, MESSAGE_ERREUR ORDER BY SOURCE, NOMBRE DESC;")
        request_error = mycursor.fetchall()

        mycursor.execute("SELECT * FROM queue_hist Where action_type in (1, 2, 3, 4, 5) AND created_date BETWEEN '"+datetime_in+"' AND '"+datetime_su+"';")
        request_ussd_success = mycursor.fetchall()

        mycursor.execute("SELECT * FROM queue Where action_type in (1, 2, 3, 4, 5) AND created_date BETWEEN '"+datetime_in+"' AND '"+datetime_su+"';")
        request_ussd_error = mycursor.fetchall()

        mydb.close()

        """ Variable email """
        # Set Global Variables
        gmail_user = 'abokor.ahmed.kadar.nour@gmail.com'
        gmail_password = 'Mail.gmail.abokor'
        # Create Email 
        mail_from = 'Abokor Ahmed-Kadar Nour'
        #mail_to = 'abdoulkader_osman@intnet.dj,abdillahi_sougueh@intnet.dj,ismael.mourad@intnet.dj,nagwa_issam@intnet.dj,mouna_zain@intnet.dj,abdek.housssein@gmail.com,dt.bss.abdifatah@gmail.com,miwsaban@gmail.com,abokor.ahmed.kadar@outlook.com'
        mail_to = 'ismael.mourad@intnet.dj,nagwa_issam@intnet.dj,mouna_zain@intnet.dj,abdek.housssein@gmail.com,dt.bss.abdifatah@gmail.com,miwsaban@gmail.com,abokor.ahmed.kadar@outlook.com'
        mail_subject = 'HEALTHCHECK TOPUP'
        mail_message_body = '''\
        Bonjour,

        Veuillez trouver ci-joint l'update des performances de la plateforme OCS/AAA/DPI et TOPUP du %s-%s (%s-%s).

        Cordialement,

        ABOKOR AHMED-KADAR NOUR
        ''' % (date_yesterday, date_today, time_in, time_su)

        """ Traitement """

        tableau_error_eti = []
        tableau_error_mes = []
        tableau_error_tot = []

        somme_error = 0

        if (len(request_error)==0):

            tableau_error_eti = 0
            tableau_error_mes = 0
            tableau_error_tot = 0

        else:

            for x in range(len(request_error)):
                tableau_error_eti.append(request_error[x][1])
                tableau_error_mes.append(request_error[x][2])
                tableau_error_tot.append(request_error[x][3])
                somme_error = somme_error + request_error[x][3]

        # initialise data of lists.
        data_error = {'Etiquettes':tableau_error_eti,
                'Message d erreur':tableau_error_mes,
                'Total':tableau_error_tot}
        
        # Create DataFrame
        t_error = pd.DataFrame(data_error)
        ###################################################################

        tableau_success_crm_sou = []
        tableau_success_crm_cod = []
        tableau_success_crm_nob = []
        tableau_success_crm_fra = []
        tableau_success_crm_tot = []

        tableau_success_ussd_sou = []
        tableau_success_ussd_cod = []
        tableau_success_ussd_nob = []
        tableau_success_ussd_fra = []
        tableau_success_ussd_tot = []

        if (len(request_error)==0):

            tableau_success_crm_sou = 0
            tableau_success_crm_cod = 0
            tableau_success_crm_nob = 0
            tableau_success_crm_fra = 0
            tableau_success_crm_tot = 0

            tableau_success_ussd_sou = 0
            tableau_success_ussd_cod = 0
            tableau_success_ussd_nob = 0
            tableau_success_ussd_fra = 0
            tableau_success_ussd_tot = 0

        else:
            tableau_success_crm_leg = []
            tableau_success_ussd_leg = []
            for x in range(len(request_success)):
                if request_success[x][0]=="CRM":
                    tableau_success_crm_leg.append(chr(x+65))
                    tableau_success_crm_cod.append(request_success[x][2])
                    tableau_success_crm_nob.append(request_success[x][3])
                    tableau_success_crm_fra.append(request_success[x][4])
                    tableau_success_crm_tot.append(request_success[x][5])
                else:
                    tableau_success_ussd_leg.append(chr(x+65))
                    tableau_success_ussd_cod.append(request_success[x][2])
                    tableau_success_ussd_nob.append(request_success[x][3])
                    tableau_success_ussd_fra.append(request_success[x][4])
                    tableau_success_ussd_tot.append(request_success[x][5])

        # initialise data of lists.
        data_success_crm = {'Etiquettes':tableau_success_crm_cod,
                'Nombre':tableau_success_crm_nob,
                'Frais':tableau_success_crm_fra,
                'Total par package':tableau_success_crm_tot}

        data_success_ussd = {'Etiquettes':tableau_success_ussd_cod,
                'Nombre':tableau_success_ussd_nob,
                'Frais':tableau_success_ussd_fra,
                'Total par package':tableau_success_ussd_tot}
        
        # Create DataFrame
        t_success_crm = pd.DataFrame(data_success_crm)

        x_crm = np.array(tableau_success_crm_cod)
        x_crm = np.array(tableau_success_crm_leg)
        y_crm = np.array(tableau_success_crm_nob)
        plt.bar(x_crm,y_crm)
        plt.savefig('./tmp/'+'success_crm.png')
        plt.close()

        t_success_ussd = pd.DataFrame(data_success_ussd)

        x_ussd = np.array(tableau_success_ussd_cod)
        x_ussd = np.array(tableau_success_ussd_leg)
        y_ussd = np.array(tableau_success_ussd_nob)
        plt.bar(x_ussd,y_ussd)
        plt.savefig('./tmp/'+'success_ussd.png')
        plt.close()

        ###################################################################

        prs = Presentation()

        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Health Check and PERFORMANCE TOPUP "
        subtitle.text = "Activity REPORT du "+date_yesterday+" – "+date_today+" ( "+time_in+" – "+time_su+" )"

        #############
        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = 'REQUEST FAILED USSD'

        rows = len(tableau_error_eti)+2
        cols = 3
        left = Inches(1.0)
        top = Inches(2.0)
        width = Inches(6.0)
        height = Inches(0.8)

        table = shapes.add_table(rows, cols, left, top, width, height).table

        # set column widths
        table.columns[0].width = Inches(2.0)
        table.columns[1].width = Inches(4.0)
        table.columns[2].width = Inches(2.0)

        # write column headings
        table.cell(0, 0).text = "Etiquettes"
        table.cell(0, 1).text = "Message d erreur"
        table.cell(0, 2).text = "Total"

        # write body cells
        k=0+1
        somme = 0
        for x in range(len(tableau_error_eti)):
            table.cell(k, 0).text = str(tableau_error_eti[x])
            table.cell(k, 1).text = str(tableau_error_mes[x])
            table.cell(k, 2).text = str(tableau_error_tot[x])
            somme = somme+tableau_error_tot[x]
            k=k+1

        table.cell(k, 0).text = "Total"
        table.cell(k, 1).text = ""
        table.cell(k, 2).text = str(somme)

        #############

        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = 'REQUEST SUCCESS FOR CRM'

        rows = len(tableau_success_crm_cod)+2
        cols = 4
        left = Inches(0.0)
        top = Inches(2.0)
        width = Inches(6.0)
        height = Inches(0.8)

        table = shapes.add_table(rows, cols, left, top, width, height).table

        # set column widths
        table.columns[0].width = Inches(4.0)
        table.columns[1].width = Inches(2.0)
        table.columns[2].width = Inches(2.0)
        table.columns[3].width = Inches(2.0)

        # write column headings
        table.cell(0, 0).text = "Etiquettes"
        table.cell(0, 1).text = "Somme de Nombre"
        table.cell(0, 2).text = "Somme de Frais"
        table.cell(0, 3).text = "Total_Per_Package"

        # write body cells
        k=0+1
        somme1 = somme2 = somme3 = 0
        for x in range(len(tableau_success_crm_cod)):
            table.cell(k, 0).text = str(tableau_success_crm_cod[x])
            table.cell(k, 1).text = str(tableau_success_crm_nob[x])
            table.cell(k, 2).text = str(tableau_success_crm_fra[x])
            table.cell(k, 3).text = str(tableau_success_crm_tot[x])
            somme1 = somme1+tableau_success_crm_nob[x]
            somme2 = somme2+tableau_success_crm_fra[x]
            somme3 = somme3+tableau_success_crm_tot[x]
            k=k+1

        table.cell(k, 0).text = "Total"
        table.cell(k, 1).text = str(somme1)
        table.cell(k, 2).text = str(somme2)
        table.cell(k, 3).text = str(somme3)

        #############

        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = 'REQUEST SUCCESS FOR CRM PART 2'

        img_path = './tmp/success_crm.png'

        left = top = Inches(2)
        pic = slide.shapes.add_picture(img_path, left, top)

        #############

        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = 'REQUEST SUCCESS FOR USSD'

        rows = len(tableau_success_ussd_cod)+2
        cols = 4
        left = Inches(0.0)
        top = Inches(2.0)
        width = Inches(6.0)
        height = Inches(0.8)

        table = shapes.add_table(rows, cols, left, top, width, height).table

        # set column widths
        table.columns[0].width = Inches(4.0)
        table.columns[1].width = Inches(2.0)
        table.columns[2].width = Inches(2.0)
        table.columns[3].width = Inches(2.0)

        # write column headings
        table.cell(0, 0).text = "Etiquettes"
        table.cell(0, 1).text = "Somme de Nombre"
        table.cell(0, 2).text = "Somme de Frais"
        table.cell(0, 3).text = "Total_Per_Package"

        # write body cells
        k=0+1
        sommeA = sommeB = sommeC = 0
        for x in range(len(tableau_success_ussd_cod)):
            table.cell(k, 0).text = str(tableau_success_ussd_cod[x])
            table.cell(k, 1).text = str(tableau_success_ussd_nob[x])
            table.cell(k, 2).text = str(tableau_success_ussd_fra[x])
            table.cell(k, 3).text = str(tableau_success_ussd_tot[x])
            sommeA = sommeA+tableau_success_ussd_nob[x]
            sommeB = sommeB+tableau_success_ussd_fra[x]
            sommeC = sommeC+tableau_success_ussd_tot[x]
            k=k+1

        table.cell(k, 0).text = "Total"
        table.cell(k, 1).text = str(sommeA)
        table.cell(k, 2).text = str(sommeB)
        table.cell(k, 3).text = str(sommeC)

        #############

        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = 'REQUEST SUCCESS FOR USSD PART 2'

        img_path = './tmp/success_ussd.png'

        left = top = Inches(2)
        pic = slide.shapes.add_picture(img_path, left, top)

        #############

        title_only_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = 'Légende pour les graphiques :'

        body_shape = shapes.placeholders[1]

        tf = body_shape.text_frame
        tf.text = 'Légende graphe :'

        p = tf.add_paragraph()
        p.text = 'CRM :'
        p.level = 1

        for x in range(len(tableau_success_crm_leg)):
            p = tf.add_paragraph()
            p.text = str(tableau_success_crm_leg[x])+" : "+str(tableau_success_crm_cod[x])
            p.level = 2

        #############

        title_only_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = 'Légende pour les graphiques :'

        body_shape = shapes.placeholders[1]

        tf = body_shape.text_frame
        tf.text = 'Légende graphe :'

        p = tf.add_paragraph()
        p.text = 'USSD :'
        p.level = 1

        for x in range(len(tableau_success_ussd_leg)):
            p = tf.add_paragraph()
            p.text = str(tableau_success_ussd_leg[x])+" : "+str(tableau_success_ussd_cod[x])
            p.level = 2

        #############

        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = 'TOTAL CRM AND USSD'

        rows = 4
        cols = 2
        left = Inches(2.0)
        top = Inches(3.0)
        width = Inches(6.0)
        height = Inches(0.8)

        table = shapes.add_table(rows, cols, left, top, width, height).table

        # set column widths
        table.columns[0].width = Inches(2.0)
        table.columns[1].width = Inches(4.0)
        # write column headings
        table.cell(0, 0).text = ''
        table.cell(0, 1).text = 'TOTAL FDJ'
        # write body cells
        table.cell(1, 0).text = 'CRM'
        table.cell(1, 1).text = str(somme3)
        table.cell(2, 0).text = 'USSD'
        table.cell(2, 1).text = str(sommeC)
        table.cell(3, 0).text = 'Total'
        total_somme = somme3+sommeC
        table.cell(3, 1).text = str(total_somme)

        #############

        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = 'Stats'

        rows = 4
        cols = 2
        left = Inches(2.0)
        top = Inches(3.0)
        width = Inches(6.0)
        height = Inches(0.8)

        table = shapes.add_table(rows, cols, left, top, width, height).table

        # set column widths
        table.columns[0].width = Inches(2.0)
        table.columns[1].width = Inches(4.0)
        # write column headings
        table.cell(0, 0).text = 'Status'
        table.cell(0, 1).text = 'Transactions'
        # write body cells
        total_somme_success = somme1+sommeA
        table.cell(1, 0).text = 'Success'
        table.cell(1, 1).text = str(total_somme_success)
        table.cell(2, 0).text = 'Failed'
        table.cell(2, 1).text = str(somme_error)
        table.cell(3, 0).text = 'Total'
        ttotal = total_somme_success + somme_error
        table.cell(3, 1).text = str(ttotal)
        #############

        pptx_name = "HealthCheck_Topup_"+date_today+".pptx"
        prs.save('./media/topup/'+pptx_name)

        ###################################################################

        # Healthcheck Ocs

        ###################################################################

        #############
        # AAA-1
        #############

        host = "10.39.234.25"
        port = 22
        username = "abokoran"
        password = "DJT_Abo@2020#"

        ssh = SSHClient()
        ssh.set_missing_host_key_policy(AutoAddPolicy())
        ssh.connect(host, port, username, password)

        command = "cd /home/scripts/HC/Logs/ && ls"
        stdin, stdout, stderr = ssh.exec_command(command)
        lines = stdout.readlines()
        log_aaa = lines[-1].strip()
        # SCPCLient takes a paramiko transport as an argument
        scp = SCPClient(ssh.get_transport())
        # upload it directly from memory
        file = '/home/scripts/HC/Logs/'+lines[-1].strip()
        scp.get(file,local_path='./tmp/')
        # close connection
        scp.close()

        command = "cd /home/scripts/HC/Output/ && ls grep Access_Result_AAA_PRD-3A-APP-01_*"
        stdin, stdout, stderr = ssh.exec_command(command)
        lines = stdout.readlines()
        file_access = lines[-1].strip()
        # SCPCLient takes a paramiko transport as an argument
        scp = SCPClient(ssh.get_transport())
        # upload it directly from memory
        file = '/home/scripts/HC/Output/'+lines[-1].strip()
        scp.get(file,local_path='./tmp/')
        # close connection
        scp.close()

        command = "cd /home/scripts/HC/Output/ && ls grep Terminate_Cause_AAA_PRD-3A-APP-01_*"
        stdin, stdout, stderr = ssh.exec_command(command)
        lines = stdout.readlines()
        file_term = lines[-1].strip()
        # SCPCLient takes a paramiko transport as an argument
        scp = SCPClient(ssh.get_transport())
        # upload it directly from memory
        file = '/home/scripts/HC/Output/'+lines[-1].strip()
        scp.get(file,local_path='./tmp/')
        # close connection
        scp.close()

        ssh.close()

        #############
        # OCS-STD
        #############

        host = "10.39.234.52"
        port = 22
        username = "abokoran"
        password = "DJT_Abo@2020#"

        ssh = SSHClient()
        ssh.set_missing_host_key_policy(AutoAddPolicy())
        ssh.connect(host, port, username, password)

        command = "cd /tmp/HC/Logs/ && ls"
        stdin, stdout, stderr = ssh.exec_command(command)
        lines = stdout.readlines()
        log_ocs = lines[-1].strip()
        # SCPCLient takes a paramiko transport as an argument
        scp = SCPClient(ssh.get_transport())
        # upload it directly from memory
        file = '/tmp/HC/Logs/'+lines[-1].strip()
        scp.get(file,local_path='./tmp/')
        # close connection
        scp.close()

        command = "cd /tmp/HC/Output/ && ls grep END_REASON_OCS_ocsstd.djiboutitelecom.dj_*"
        stdin, stdout, stderr = ssh.exec_command(command)
        lines = stdout.readlines()
        file_end = lines[-1].strip()
        # SCPCLient takes a paramiko transport as an argument
        scp = SCPClient(ssh.get_transport())
        # upload it directly from memory
        file = '/tmp/HC/Output/'+lines[-1].strip()
        scp.get(file,local_path='./tmp/')
        # close connection
        scp.close()

        ssh.close()

        #########################################################################

        """ La date et l'heure """
        time_in = "09:00:00"
        time_su = "09:00:00"
        datetime_today = datetime.now()
        date_today = datetime_today.strftime("%Y-%m-%d")
        datetime_yesterday = datetime.now() - timedelta(1)
        date_yesterday = datetime_yesterday.strftime("%Y-%m-%d")
        date_yesterday_ex = datetime_yesterday.strftime("%Y%m%d")
        datetime_in = date_yesterday+" "+time_in
        datetime_su = date_today+" "+time_su

        timestamp_doday = int(datetime.timestamp(datetime_today))
        timestamp_yesterday = int(datetime.timestamp(datetime_yesterday))

        #########################################################################

        prs = Presentation()

        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Health Check and PERFORMANCE OCS/AAA/DPI "
        subtitle.text = "Activity REPORT du "+date_yesterday+" – "+date_today+" ( "+time_in+" – "+time_su+" )"

        #########################################################################

        r=requests.post('http://10.39.234.54/centreon/index.php')

        # PHPSESSID
        cooki = r.cookies['PHPSESSID']
        cookiess = dict(PHPSESSID=cooki)

        # Centreon_token
        for item in r.text.split("\n"):
            if "centreon_token" in item:
                ligne = item.strip()
        value = ligne[50:82]

        # Paramètre
        payload = {'useralias': 'dt_supervision', 'password': 'DT2020$', 'submitLogin': 'Connect', 'centreon_token': value}

        # Connexion
        r=requests.post('http://10.39.234.54/centreon/index.php',data=payload,cookies=cookiess)

        #########################################################################

        lists_services = [
            '100_312', #Caps
            '100_329', #Diameter Open data
            '100_324', #Diameter Active sessions
            '100_325', #Diameter Call active
            '100_327', #Diameter data counters
            '100_279', #Memory Usage Ocs-std
            '99_227', #Memory Usage Ocs-Mng
            '100_265', #Int_eth0 ocs-std
            '100_266', #Int_eth1 ocs-std
            '100_267', #Int_eth2 ocs-std
        ]

        list_graph = []

        for x in lists_services:
            lien = 'http://10.39.234.54/centreon/api/internal.php?object=centreon_metric&action=metricsDataByService&ids='+x+'&start='+str(timestamp_yesterday)+'&end='+str(timestamp_doday)
            s = requests.get(lien,cookies=cookiess)
            list_graph.append(s.text)

        # Déconnexion
        deco = {'disconnect': 1}
        t = requests.get('http://10.39.234.54/centreon/index.php',params=deco,cookies=cookiess)

        # Extraction des abs
        list_abs = []
        extraction = list_graph[0]

        i=0
        for x in extraction:
            
            if x=='[':
                i=i+1

            if i==5 and x!='[' and x!=']':
                list_abs.append(x)

            if i==5 and x==']':
                break

        i=0
        result = ''
        abs = []

        for x in range(len(list_abs)):
            
            if list_abs[x]!=',':
                result = result+list_abs[x]

            if list_abs[x] == ',':
                abs.append(result)
                result = ''

            if x == len(list_abs)-1:
                abs.append(result)
                result=''

        for x in range(len(abs)):
            abs[x]=datetime.fromtimestamp(int(abs[x]))

        # Extraction des ord Caps

        list_data = []

        i=0
        list_ord = []

        for y in list_graph[0]:
            
            if y=='[':
                i=i+1

            if i==3 and y!='[' and y!=']':
                list_ord.append(y)

            if i==3 and y==']':
                break

        list_data.append(list_ord)

        i=0
        list_ord = []

        for y in list_graph[0]:
            
            if y=='[':
                i=i+1

            if i==4 and y!='[' and y!=']':
                list_ord.append(y)

            if i==4 and y==']':
                break

        list_data.append(list_ord)
        ################################################
        data_string = []

        for x in range(len(list_data)):

            ordd = []
            result = ''

            for y in range(len(list_data[x])):

                if list_data[x][y]!=',':
                    result = result+list_data[x][y]

                if list_data[x][y] == ',':
                    ordd.append(result)
                    result = ''

                if y == len(list_data[x])-1:
                    ordd.append(result)
                    result=''

            data_string.append(ordd)

        ###############################################################

        data = []

        for x in data_string:
            
            float_data = []

            for y in x:

                if y == 'null':
                    float_data.append(float(0))
                else:
                    float_data.append(float(y))

            data.append(float_data)

        ###############################################################

        fig = go.Figure()
        fig.add_trace(go.Scatter(x=abs, y=data[0], name='refused',line = dict(color='rgba(255, 99, 71, 0.8)')))
        fig.add_trace(go.Scatter(x=abs, y=data[1],name='caps', line = dict(color='rgba(135, 211, 124, 1)'),fill='tozeroy'))
        fig.update_layout(hovermode="x unified")                                
        fig.write_image("./tmp/caps.png",width=650, height=400, scale=1)

        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = 'Caps'

        img_path = './tmp/caps.png'

        left = Inches(0.5)
        top = Inches(1.5)
        pic = slide.shapes.add_picture(img_path, left, top)
        ######

        # Extraction des Diameter Open data

        list_data = []

        i=0
        list_ord = []

        for y in list_graph[1]:
            
            if y=='[':
                i=i+1

            if i==3 and y!='[' and y!=']':
                list_ord.append(y)

            if i==3 and y==']':
                break

        list_data.append(list_ord)
        ################################################
        data_string = []

        for x in range(len(list_data)):

            ordd = []
            result = ''

            for y in range(len(list_data[x])):

                if list_data[x][y]!=',':
                    result = result+list_data[x][y]

                if list_data[x][y] == ',':
                    ordd.append(result)
                    result = ''

                if y == len(list_data[x])-1:
                    ordd.append(result)
                    result=''

            data_string.append(ordd)

        ###############################################################

        data = []

        for x in data_string:
            
            float_data = []

            for y in x:

                if y == 'null':
                    float_data.append(float(0))
                else:
                    float_data.append(float(y))

            data.append(float_data)

        ###############################################################

        fig = go.Figure()
        fig.add_trace(go.Scatter(x=abs, y=data[0], name='Application_DIAMETER_Data_open_data_sessions',line = dict(color='rgba(255, 99, 71, 0.8)')))
        fig.update_layout(hovermode="x unified")                                
        fig.write_image("./tmp/diameter_data_open_data_sessions.png",width=650, height=400, scale=1)

        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = 'DIAMETER Data open data sessions'

        img_path = './tmp/diameter_data_open_data_sessions.png'

        left = Inches(0.5)
        top = Inches(1.5)
        pic = slide.shapes.add_picture(img_path, left, top)
        ######

        # Extraction des Diameter Open data

        list_data = []

        i=0
        list_ord = []

        for y in list_graph[2]:
            
            if y=='[':
                i=i+1

            if i==3 and y!='[' and y!=']':
                list_ord.append(y)

            if i==3 and y==']':
                break

        list_data.append(list_ord)
        ################################################
        data_string = []

        for x in range(len(list_data)):

            ordd = []
            result = ''

            for y in range(len(list_data[x])):

                if list_data[x][y]!=',':
                    result = result+list_data[x][y]

                if list_data[x][y] == ',':
                    ordd.append(result)
                    result = ''

                if y == len(list_data[x])-1:
                    ordd.append(result)
                    result=''

            data_string.append(ordd)

        ###############################################################

        data = []

        for x in data_string:
            
            float_data = []

            for y in x:

                if y == 'null':
                    float_data.append(float(0))
                else:
                    float_data.append(float(y))

            data.append(float_data)

        ###############################################################

        fig = go.Figure()
        fig.add_trace(go.Scatter(x=abs, y=data[0], name='DIAMETER_server_sessions_active',line = dict(color='rgba(84, 176, 255, 0.8)')))
        fig.update_layout(hovermode="x unified")                                
        fig.write_image("./tmp/diameter_server_sessions_active.png",width=650, height=400, scale=1)

        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = 'Diameter server sessions active'

        img_path = './tmp/diameter_server_sessions_active.png'

        left = Inches(0.5)
        top = Inches(1.5)
        pic = slide.shapes.add_picture(img_path, left, top)
        ######

        # Extraction des Diameter call active

        list_data = []

        i=0
        list_ord = []

        for y in list_graph[3]:
            
            if y=='[':
                i=i+1

            if i==3 and y!='[' and y!=']':
                list_ord.append(y)

            if i==3 and y==']':
                break

        list_data.append(list_ord)

        i=0
        list_ord = []

        for y in list_graph[3]:
            
            if y=='[':
                i=i+1

            if i==4 and y!='[' and y!=']':
                list_ord.append(y)

            if i==4 and y==']':
                break

        list_data.append(list_ord)
        ################################################
        data_string = []

        for x in range(len(list_data)):

            ordd = []
            result = ''

            for y in range(len(list_data[x])):

                if list_data[x][y]!=',':
                    result = result+list_data[x][y]

                if list_data[x][y] == ',':
                    ordd.append(result)
                    result = ''

                if y == len(list_data[x])-1:
                    ordd.append(result)
                    result=''

            data_string.append(ordd)

        ###############################################################

        data = []

        for x in data_string:
            
            float_data = []

            for y in x:

                if y == 'null':
                    float_data.append(float(0))
                else:
                    float_data.append(float(y))

            data.append(float_data)

        ###############################################################

        fig = go.Figure()
        fig.add_trace(go.Scatter(x=abs, y=data[0], name='MO_calls_active',line = dict(color='rgba(48, 255, 152, 0.8)')))
        fig.add_trace(go.Scatter(x=abs, y=data[1],name='MT_calls_active', line = dict(color='rgba(29, 148, 88, 1)')))
        fig.update_layout(hovermode="x unified")                                
        fig.write_image("./tmp/diameter_calls_active.png",width=650, height=400, scale=1)

        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = 'Diameter calls active'

        img_path = './tmp/diameter_calls_active.png'

        left = Inches(0.5)
        top = Inches(1.5)
        pic = slide.shapes.add_picture(img_path, left, top)
        ######

        # Extraction des Diameter data counter

        list_data = []

        i=0
        list_ord = []

        for y in list_graph[4]:
            
            if y=='[':
                i=i+1

            if i==3 and y!='[' and y!=']':
                list_ord.append(y)

            if i==3 and y==']':
                break

        list_data.append(list_ord)

        i=0
        list_ord = []

        for y in list_graph[4]:
            
            if y=='[':
                i=i+1

            if i==4 and y!='[' and y!=']':
                list_ord.append(y)

            if i==4 and y==']':
                break

        list_data.append(list_ord)

        i=0
        list_ord = []

        for y in list_graph[4]:
            
            if y=='[':
                i=i+1

            if i==5 and y!='[' and y!=']':
                list_ord.append(y)

            if i==5 and y==']':
                break

        list_data.append(list_ord)
        ################################################
        data_string = []

        for x in range(len(list_data)):

            ordd = []
            result = ''

            for y in range(len(list_data[x])):

                if list_data[x][y]!=',':
                    result = result+list_data[x][y]

                if list_data[x][y] == ',':
                    ordd.append(result)
                    result = ''

                if y == len(list_data[x])-1:
                    ordd.append(result)
                    result=''

            data_string.append(ordd)

        ###############################################################

        data = []

        for x in data_string:
            
            float_data = []

            for y in x:

                if y == 'null':
                    float_data.append(float(0))
                else:
                    float_data.append(float(y))

            data.append(float_data)

        ###############################################################

        fig = go.Figure()
        fig.add_trace(go.Scatter(x=abs, y=data[0], name='closed_data_sessions',line = dict(color='rgba(159, 255, 208, 0.8)')))
        fig.add_trace(go.Scatter(x=abs, y=data[1],name='sessions_rejected_network', line = dict(color='rgba(229, 255, 92, 1)')))
        fig.add_trace(go.Scatter(x=abs, y=data[2],name='sessions_rejected_ocs', line = dict(color='rgba(116, 161, 196, 1)')))
        fig.update_layout(hovermode="x unified")                                
        fig.write_image("./tmp/diameter_data_counter.png",width=650, height=400, scale=1)

        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = 'Diameter data counter'

        img_path = './tmp/diameter_data_counter.png'

        left = Inches(0.5)
        top = Inches(1.5)
        pic = slide.shapes.add_picture(img_path, left, top)
        ######

        # Extraction des Memo std

        list_data = []

        i=0
        list_ord = []

        for y in list_graph[5]:
            
            if y=='[':
                i=i+1

            if i==3 and y!='[' and y!=']':
                list_ord.append(y)

            if i==3 and y==']':
                break

        list_data.append(list_ord)

        i=0
        list_ord = []

        for y in list_graph[5]:
            
            if y=='[':
                i=i+1

            if i==4 and y!='[' and y!=']':
                list_ord.append(y)

            if i==4 and y==']':
                break

        list_data.append(list_ord)
        ################################################
        data_string = []

        for x in range(len(list_data)):

            ordd = []
            result = ''

            for y in range(len(list_data[x])):

                if list_data[x][y]!=',':
                    result = result+list_data[x][y]

                if list_data[x][y] == ',':
                    ordd.append(result)
                    result = ''

                if y == len(list_data[x])-1:
                    ordd.append(result)
                    result=''

            data_string.append(ordd)

        ###############################################################

        data = []

        for x in data_string:
            
            float_data = []

            for y in x:

                if y == 'null':
                    float_data.append(float(0))
                else:
                    float_data.append(float(y))

            data.append(float_data)

        ###############################################################

        fig = go.Figure()
        fig.add_trace(go.Scatter(x=abs, y=data[0], name='Used memory',line = dict(color='rgba(135, 211, 124, 1)'),fill='tozeroy'))
        fig.add_trace(go.Scatter(x=abs, y=data[1],name='Used swap', line = dict(color='rgba(100, 24, 129, 1)'),fill='tozeroy'))
        fig.update_layout(hovermode="x unified")                                
        fig.write_image("./tmp/memo_std.png",width=650, height=400, scale=1)

        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = 'Memory_Usage ocs_std'

        img_path = './tmp/memo_std.png'

        left = Inches(0.5)
        top = Inches(1.5)
        pic = slide.shapes.add_picture(img_path, left, top)
        ######

        # Extraction des Memo mng

        list_data = []

        i=0
        list_ord = []

        for y in list_graph[6]:
            
            if y=='[':
                i=i+1

            if i==3 and y!='[' and y!=']':
                list_ord.append(y)

            if i==3 and y==']':
                break

        list_data.append(list_ord)

        i=0
        list_ord = []

        for y in list_graph[6]:
            
            if y=='[':
                i=i+1

            if i==4 and y!='[' and y!=']':
                list_ord.append(y)

            if i==4 and y==']':
                break

        list_data.append(list_ord)
        ################################################
        data_string = []

        for x in range(len(list_data)):

            ordd = []
            result = ''

            for y in range(len(list_data[x])):

                if list_data[x][y]!=',':
                    result = result+list_data[x][y]

                if list_data[x][y] == ',':
                    ordd.append(result)
                    result = ''

                if y == len(list_data[x])-1:
                    ordd.append(result)
                    result=''

            data_string.append(ordd)

        ###############################################################

        data = []

        for x in data_string:
            
            float_data = []

            for y in x:

                if y == 'null':
                    float_data.append(float(0))
                else:
                    float_data.append(float(y))

            data.append(float_data)

        ###############################################################

        fig = go.Figure()
        fig.add_trace(go.Scatter(x=abs, y=data[0], name='Used memory',line = dict(color='rgba(135, 211, 124, 1)'),fill='tozeroy'))
        fig.add_trace(go.Scatter(x=abs, y=data[1],name='Used swap', line = dict(color='rgba(100, 24, 129, 1)'),fill='tozeroy'))
        fig.update_layout(hovermode="x unified")                                
        fig.write_image("./tmp/memo_mng.png",width=650, height=400, scale=1)

        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = 'Memory_Usage ocs_mng'

        img_path = './tmp/memo_mng.png'

        left = Inches(0.5)
        top = Inches(1.5)
        pic = slide.shapes.add_picture(img_path, left, top)
        ######

        # Extraction des Ocsstd_eth0

        list_data = []

        i=0
        list_ord = []

        for y in list_graph[7]:
            
            if y=='[':
                i=i+1

            if i==3 and y!='[' and y!=']':
                list_ord.append(y)

            if i==3 and y==']':
                break

        list_data.append(list_ord)

        i=0
        list_ord = []

        for y in list_graph[7]:
            
            if y=='[':
                i=i+1

            if i==4 and y!='[' and y!=']':
                list_ord.append(y)

            if i==4 and y==']':
                break

        list_data.append(list_ord)
        ################################################
        data_string = []

        for x in range(len(list_data)):

            ordd = []
            result = ''

            for y in range(len(list_data[x])):

                if list_data[x][y]!=',':
                    result = result+list_data[x][y]

                if list_data[x][y] == ',':
                    ordd.append(result)
                    result = ''

                if y == len(list_data[x])-1:
                    ordd.append(result)
                    result=''

            data_string.append(ordd)

        ###############################################################

        data = []

        for x in data_string:
            
            float_data = []

            for y in x:

                if y == 'null':
                    float_data.append(float(0))
                else:
                    float_data.append(float(y))

            data.append(float_data)

        ###############################################################

        fig = go.Figure()
        fig.add_trace(go.Scatter(x=abs, y=data[0],name='traffic_out',line = dict(color='rgba(210, 78, 100, 1)'),fill='tozeroy'))
        fig.add_trace(go.Scatter(x=abs, y=data[1],name='traffic_in', line = dict(color='rgba(135, 211, 124, 1)'),fill='tozeroy'))
        fig.update_layout(hovermode="x unified")                                
        fig.write_image("./tmp/eth0.png",width=650, height=400, scale=1)

        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = 'Ocs_std Eth0'

        img_path = './tmp/eth0.png'

        left = Inches(0.5)
        top = Inches(1.5)
        pic = slide.shapes.add_picture(img_path, left, top)
        ######

        # Extraction des Ocsstd_eth1

        list_data = []

        i=0
        list_ord = []

        for y in list_graph[8]:
            
            if y=='[':
                i=i+1

            if i==3 and y!='[' and y!=']':
                list_ord.append(y)

            if i==3 and y==']':
                break

        list_data.append(list_ord)

        i=0
        list_ord = []

        for y in list_graph[8]:
            
            if y=='[':
                i=i+1

            if i==4 and y!='[' and y!=']':
                list_ord.append(y)

            if i==4 and y==']':
                break

        list_data.append(list_ord)
        ################################################
        data_string = []

        for x in range(len(list_data)):

            ordd = []
            result = ''

            for y in range(len(list_data[x])):

                if list_data[x][y]!=',':
                    result = result+list_data[x][y]

                if list_data[x][y] == ',':
                    ordd.append(result)
                    result = ''

                if y == len(list_data[x])-1:
                    ordd.append(result)
                    result=''

            data_string.append(ordd)

        ###############################################################

        data = []

        for x in data_string:
            
            float_data = []

            for y in x:

                if y == 'null':
                    float_data.append(float(0))
                else:
                    float_data.append(float(y))

            data.append(float_data)

        ###############################################################

        fig = go.Figure()
        fig.add_trace(go.Scatter(x=abs, y=data[0],name='traffic_out',line = dict(color='rgba(210, 78, 100, 1)'),fill='tozeroy'))
        fig.add_trace(go.Scatter(x=abs, y=data[1],name='traffic_in', line = dict(color='rgba(135, 211, 124, 1)'),fill='tozeroy'))
        fig.update_layout(hovermode="x unified")                                
        fig.write_image("./tmp/eth1.png",width=650, height=400, scale=1)

        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = 'Ocs_std Eth1'

        img_path = './tmp/eth1.png'

        left = Inches(0.5)
        top = Inches(1.5)
        pic = slide.shapes.add_picture(img_path, left, top)
        ######

        # Extraction des Ocsstd_eth2

        list_data = []

        i=0
        list_ord = []

        for y in list_graph[9]:
            
            if y=='[':
                i=i+1

            if i==3 and y!='[' and y!=']':
                list_ord.append(y)

            if i==3 and y==']':
                break

        list_data.append(list_ord)

        i=0
        list_ord = []

        for y in list_graph[9]:
            
            if y=='[':
                i=i+1

            if i==4 and y!='[' and y!=']':
                list_ord.append(y)

            if i==4 and y==']':
                break

        list_data.append(list_ord)
        ################################################
        data_string = []

        for x in range(len(list_data)):

            ordd = []
            result = ''

            for y in range(len(list_data[x])):

                if list_data[x][y]!=',':
                    result = result+list_data[x][y]

                if list_data[x][y] == ',':
                    ordd.append(result)
                    result = ''

                if y == len(list_data[x])-1:
                    ordd.append(result)
                    result=''

            data_string.append(ordd)

        ###############################################################

        data = []

        for x in data_string:
            
            float_data = []

            for y in x:

                if y == 'null':
                    float_data.append(float(0))
                else:
                    float_data.append(float(y))

            data.append(float_data)

        ###############################################################

        fig = go.Figure()
        fig.add_trace(go.Scatter(x=abs, y=data[0],name='traffic_out',line = dict(color='rgba(210, 78, 100, 1)'),fill='tozeroy'))
        fig.add_trace(go.Scatter(x=abs, y=data[1],name='traffic_in', line = dict(color='rgba(135, 211, 124, 1)'),fill='tozeroy'))
        fig.update_layout(hovermode="x unified")                                
        fig.write_image("./tmp/eth2.png",width=650, height=400, scale=1)

        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = 'Ocs_std Eth2'

        img_path = './tmp/eth2.png'

        left = Inches(0.5)
        top = Inches(1.5)
        pic = slide.shapes.add_picture(img_path, left, top)

        ###############################################################################################
        f=r"tmp//"+file_access
        col_list = ["COUNT","DATE_HOUR","RETURN_CODE"]
        QueryResult = pd.read_csv(f, usecols=col_list, sep=';')

        table = pd.pivot_table(QueryResult, values='COUNT', index=['DATE_HOUR'],
                            columns=['RETURN_CODE'], aggfunc=np.sum)

        table.plot()
        plt.xticks(rotation=90)
        plt.tight_layout()
        plt.savefig('./tmp/access_result_aaa.png')
        plt.close()

        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = 'Access Result AAA'

        img_path = './tmp/access_result_aaa.png'

        left = Inches(0.5)
        top = Inches(1.5)
        pic = slide.shapes.add_picture(img_path, left, top)

        #####

        f=r"tmp//"+file_term
        col_list1 = ["Count","Acct-Terminate-Cause","Event-Timestamp-Date"]
        QueryResult1 = pd.read_csv(f, usecols=col_list1, sep=';')
        list = []
        for x in QueryResult1["Event-Timestamp-Date"]:
            annee = str(x)[0:4]
            mois = str(x)[4:6]
            jour = str(x)[6:8]
            date = annee+"-"+mois+"-"+jour
            list.append(date)

        for x in range(len(QueryResult1["Event-Timestamp-Date"])):
            QueryResult1["Event-Timestamp-Date"][x]=list[x]
            
        table1 = pd.pivot_table(QueryResult1, values='Count', index=['Event-Timestamp-Date'],
                            columns=['Acct-Terminate-Cause'], aggfunc=np.sum)
        table1.plot()
        plt.xticks(rotation=90)
        plt.tight_layout()
        plt.savefig('./tmp/term_cause_aaa.png')
        plt.close()

        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = 'Terminate Cause AAA'

        img_path = './tmp/term_cause_aaa.png'

        left = Inches(0.5)
        top = Inches(1.5)
        pic = slide.shapes.add_picture(img_path, left, top)

        #####

        f=r"tmp//"+file_end

        col_list2 = ["COUNT","END_DATE_HOUR","END_REASON_CODE_TXT","END_DATE"]
        QueryResult2 = pd.read_csv(f, usecols=col_list2, sep=';')

        list = []
        for x in QueryResult2["END_DATE_HOUR"]:
            annee = str(x)[0:4]
            mois = str(x)[4:6]
            jour = str(x)[6:8]
            hour = str(x)[8:10]
            date = annee+"-"+mois+"-"+jour+"-"+hour+"h"
            list.append(date)

        for x in range(len(list)):
            QueryResult2["END_DATE_HOUR"][x]=list[x]

        with pd.ExcelWriter(r"tmp/New.xlsx") as writer:
            QueryResult2.to_excel(writer, sheet_name='Total')

        f=r"tmp//"+"New.xlsx"
        QueryResult3 = pd.read_excel(f)
        QueryResult3.drop(QueryResult2[QueryResult2['END_DATE'] <int(date_yesterday_ex)].index, inplace = True)
        table2 = pd.pivot_table(QueryResult3, values='COUNT', index=['END_DATE_HOUR'],
                            columns=['END_REASON_CODE_TXT'], aggfunc=np.sum)
        plt.plot(table2)
        plt.xticks(rotation=90)
        plt.tight_layout()
        plt.savefig('./tmp/end_reason_ocs.png')
        plt.close()

        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = 'END REASON OCS'

        img_path = './tmp/end_reason_ocs.png'

        left = Inches(0.5)
        top = Inches(1.5)
        pic = slide.shapes.add_picture(img_path, left, top)

        ######

        ocs_pptx_name = "HealthCheck_OCS-AAA-DPI_"+date_today+".pptx"
        prs.save('./media/ocs-aaa-dpi/'+ocs_pptx_name)

        ###################################################################

        message = MIMEMultipart('alternative')
        message['Subject'] = mail_subject
        message['From'] = mail_from
        message['To'] = mail_to
        pptx = MIMEApplication(open('./media/topup/'+pptx_name,'rb').read())
        pptx.add_header('content-disposition', 'attachment',filename=pptx_name)
        message.attach(pptx)
        pptx_ocs = MIMEApplication(open('./media/ocs-aaa-dpi/'+ocs_pptx_name,'rb').read())
        pptx_ocs.add_header('content-disposition', 'attachment',filename=ocs_pptx_name)
        message.attach(pptx_ocs)
        message.attach(MIMEText(mail_message_body, 'plain'))
        # Sent Email
        server = smtplib.SMTP_SSL('smtp.gmail.com', 465)
        server.login(gmail_user, gmail_password)
        server.sendmail(mail_from,mail_to.split(","), message.as_string())
        server.close()

        ###################################################################

        ficher = HEALTHCHECK()
        ficher.topup.name = pptx_name
        ficher.ocs_aaa_dpi.name = ocs_pptx_name
        ficher.save()

        pass