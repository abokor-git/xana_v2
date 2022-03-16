from django_extensions.management.jobs import DailyJob

import pandas as pd

import matplotlib.pyplot as plt
import numpy as np

from pptx import Presentation
from pptx.util import Inches

from datetime import datetime, timedelta

from apps.app.models import Logs, PackageData
from apps.app.models import Queue, QueueHist

import smtplib
from email.mime.text import MIMEText
from email.mime.application import MIMEApplication
from email.mime.multipart import MIMEMultipart

class Job(DailyJob):
    help = "Stats Daily TopUp at 08:00 am !"

    def execute(self):
        
        PROCESS_NUMBER = 1615

        """ La date et l'heure """
        time_in = "00:00:00"
        time_su = "00:00:00"
        datetime_today = datetime.now()
        datetime_yesterday = datetime.now()-timedelta(1)
        date_today = datetime_today.strftime("%Y-%m-%d")
        date_yesterday = datetime_yesterday.strftime("%Y-%m-%d")
        datetime_now = datetime_today.strftime("%Y-%m-%d %H:%M:%S")
        datetime_min = date_yesterday+" "+time_in
        datetime_max = date_today+" "+time_su

        ###################################################################################

        prs = Presentation()

        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Statistique TopUp Requests"
        subtitle.text = date_yesterday+" – "+date_today+" ( "+time_in+" – "+time_su+" )"

        ###################################################################################
        
        """ List of packages """
        etudiant_list = ['ETUDIANT_MOIS', 'ETUDIANT_SEMESTRE', 'ETUDIANT_ANNEE']
        etudiant_region_list = ['ETUDIANT_MOIS_REGION', 'ETUDIANT_SEMEST_REGION', 'ETUDIANT_ANNEE_REGION']

        residentiel_list = ['RESIDENTIEL_3M_MOIS', 'RESIDENTIEL_3M_SEMESTR', 'RESIDENTIEL_3M_ANNEE',
            'RESIDENTIEL_4M_MOIS', 'RESIDENTIEL_4M_SEMESTR', 'RESIDENTIEL_4M_ANNEE',
            'RESIDENTIEL_5M_MOIS', 'RESIDENTIEL_5M_SEMESTR', 'RESIDENTIEL_5M_ANNEE',
            'RESIDENTIEL_6M_MOIS', 'RESIDENTIEL_6M_SEMESTR', 'RESIDENTIEL_6M_ANNEE']
        residentiel_region_list = ['RES_3M_MOIS_REGION', 'RES_3M_SEMESTR_REGION', 'RES_3M_ANNEE_REGION',
            'RES_4M_MOIS_REGION', 'RES_4M_SEMESTR_REGION', 'RES_4M_ANNEE_REGION',
            'RES_5M_MOIS_REGION', 'RES_5M_SEMESTR_REGION', 'RES_5M_ANNEE_REGION',
            'RES_6M_MOIS_REGION', 'RES_6M_SEMESTR_REGION', 'RES_6M_ANNEE_REGION']

        gpp_list = ['GRAND_PUBLIC_1M_MOIS', 'GRAND_PUBLIC_1M_SEMEST', 'GRAND_PUBLIC_1M_ANNEE',
            'GRAND_PUBLIC_3M_MOIS', 'GRAND_PUBLIC_3M_SEMEST', 'GRAND_PUBLIC_3M_ANNEE',
            'GRAND_PUBLIC_4M_MOIS', 'GRAND_PUBLIC_4M_SEMEST', 'GRAND_PUBLIC_4M_ANNEE',
            'GRAND_PUBLIC_5M_MOIS', 'GRAND_PUBLIC_5M_SEMEST', 'GRAND_PUBLIC_5M_ANNEE',
            'GRAND_PUBLIC_6M_MOIS', 'GRAND_PUBLIC_6M_SEMEST', 'GRAND_PUBLIC_6M_ANNEE']
        gpp_region_list = ['GPP_1M_MOIS_REGION', 'GPP_1M_SEMEST_REGION', 'GPP_1M_ANNEE_REGION',
            'GPP_3M_MOIS_REGION', 'GPP_3M_SEMESTR_REGION', 'GPP_3M_ANNEE_REGION',
            'GPP_4M_MOIS_REGION', 'GPP_4M_SEMESTR_REGION', 'GPP_4M_ANNEE_REGION',
            'GPP_5M_MOIS_REGION', 'GPP_5M_SEMESTR_REGION', 'GPP_5M_ANNEE_REGION',
            'GPP_6M_MOIS_REGION', 'GPP_6M_SEMESTR_REGION', 'GPP_6M_ANNEE_REGION']

        auto_gp_list = ['AUTO_GRAND_PUBLIC_1M_MOIS',
            'AUTO_GRAND_PUBLIC_3M_MOIS','AUTO_GRAND_PUBLIC_4M_MOIS',
            'AUTO_GRAND_PUBLIC_5M_MOIS','AUTO_GRAND_PUBLIC_6M_MOIS',]

        employe_list = ['EMPLOYE_2M', 'EMPLOYE_4M', 'EMPLOYE_6M', 'EMPLOYE_8M', 'EMPLOYE_10M']

        auto_recharge_list = ['AUTO_RECHARGE_I']

        frais_retablissement = ['FRAIS_RETABLISSEMENT']

        voix_list = ['VOIX_MOIS', 'VOIX_SEMESTRE', 'VOIX_ANNEE']
        voix_region_list = ['VOIX_MOIS_REGION', 'VOIX_SEMESTRE_REGION', 'VOIX_ANNEE_REGION']

        gigasup_list = ['GIGASUP_10GO', 'GIGASUP_30GO', 'GIGASUP_50GO', 'GIGASUP_100GO', 'GIGASUP_200GO']
        voixsup_list = ['VOIXSUP_60MINS', 'VOIXSUP_150MINS', 'VOIXSUP_330MINS']

        recharge_list = ['RECHARGE_A', 'RECHARGE_B', 'RECHARGE_C', 'RECHARGE_D', 'RECHARGE_E', 'RECHARGE_EXDG', 'RECHARGE_F', 'RECHARGE_G', 'RECHARGE_H']

        ######################################################################

        list_name = []
        list_name.append("Les packages Etudiant")
        list_name.append("Les packages Etudiant Région")
        list_name.append("Les packages Résidentiel")
        list_name.append("Les packages Résidentiel Région")
        list_name.append("Les packages Grand Public Postpaid")
        list_name.append("Les packages Grand Public Postpaid Région")
        list_name.append("Les packages Auto Grand Public")
        list_name.append("Les packages Employee")
        list_name.append("Packages Auto-Recharge-I")
        list_name.append("Packages Frais de Rétablissement")
        list_name.append("Les packages Voix")
        list_name.append("Les packages Voix Région")
        list_name.append("Les packages Giga Sup")
        list_name.append("Les packages Voix Sup")
        list_name.append("Les packages Recharge Employer")
        
        list_all = []
        list_all.append(etudiant_list)
        list_all.append(etudiant_region_list)
        list_all.append(residentiel_list)
        list_all.append(residentiel_region_list)
        list_all.append(gpp_list)
        list_all.append(gpp_region_list)
        list_all.append(auto_gp_list)
        list_all.append(employe_list)
        list_all.append(auto_recharge_list)
        list_all.append(frais_retablissement)
        list_all.append(voix_list)
        list_all.append(voix_region_list)
        list_all.append(gigasup_list)
        list_all.append(voixsup_list)
        list_all.append(recharge_list)

        n=0

        for y in list_all: 
            ###################
            """ traitement """
            frais = []
            crm = []
            ussd = []
            total = []
            cout = []

            for x in y:
                nb_frais = PackageData.objects.get(code=x).frais
                nb_crm = QueueHist.objects.filter(created_date__range=[datetime_min,datetime_max],command_content__endswith=x,action_type=44).count()
                nb_ussd = QueueHist.objects.filter(created_date__range=[datetime_min,datetime_max],command_content__endswith=x,action_type=4).count()
                nb_total = nb_crm+nb_ussd
                nb_cout = nb_total * nb_frais
                
                frais.append(nb_frais)
                crm.append(nb_crm)
                ussd.append(nb_ussd)
                total.append(nb_total)
                cout.append(nb_cout)

            data = {'Package':y,'Frais':frais,'Crm':crm,'Ussd':ussd,'Total':total,'Coût':cout}
            df = pd.DataFrame(data)

            ##########################################

            if df['Total'].sum() != 0:

                title_only_slide_layout = prs.slide_layouts[5]
                slide = prs.slides.add_slide(title_only_slide_layout)
                shapes = slide.shapes

                shapes.title.text = list_name[n]

                rows = len(y)+1
                cols = 6
                left = Inches(0.2)
                top = Inches(1.5)
                width = Inches(6.0)
                height = Inches(0.8)

                table = shapes.add_table(rows, cols, left, top, width, height).table

                # set column widths
                table.columns[0].width = Inches(3.0)
                table.columns[1].width = Inches(2.0)
                table.columns[2].width = Inches(1.0)
                table.columns[3].width = Inches(1.0)
                table.columns[4].width = Inches(1.0)
                table.columns[5].width = Inches(1.5)

                # write column headings
                table.cell(0, 0).text = "Package"
                table.cell(0, 1).text = "Frais"
                table.cell(0, 2).text = "Crm"
                table.cell(0, 3).text = "Ussd"
                table.cell(0, 4).text = "Total"
                table.cell(0, 5).text = "Coût"

                # write body cells
                k=0+1
                for l in df.index:
                    table.cell(k, 0).text = str(df["Package"][l])
                    table.cell(k, 1).text = str(df["Frais"][l])
                    table.cell(k, 2).text = str(df["Crm"][l])
                    table.cell(k, 3).text = str(df["Ussd"][l])
                    table.cell(k, 4).text = str(df["Total"][l])
                    table.cell(k, 5).text = str(df["Coût"][l])
                    k=k+1

                print(df)

                fig, ax = plt.subplots()
                ax.pie(df['Coût'], labels=df['Package'], shadow=True, startangle=90)
                ax.axis('equal')
                plt.legend(title="Packages :")
                file_name = "./tmp/pie.png"
                plt.savefig(file_name)

                #########################################

                title_only_slide_layout = prs.slide_layouts[5]
                slide = prs.slides.add_slide(title_only_slide_layout)
                shapes = slide.shapes

                shapes.title.text = "Graph"

                img_path = './tmp/pie.png'

                left = top = Inches(2)
                pic = slide.shapes.add_picture(img_path, left, top)

            n=n+1

        #################################################################################

        name = "stats_topup_daily_"+date_yesterday+"_"+date_today+".pptx"
        prs.save('./stats/daily/'+name)

        #################################################################################

        link = "http://172.16.1.89/"

        """ Variable email """
        # Set Global Variables
        gmail_user = 'xana.v2.system@gmail.com'
        gmail_password = 'Xana_v2_system'
        # Create Email 
        mail_from = 'Xana System'
        #mail_to = 'nagwa_issam@intnet.dj,mouna_zain@intnet.dj,abdek.housssein@gmail.com,dt.bss.abdifatah@gmail.com,abokor.ahmed.kadar@outlook.com'
        mail_to = 'abokor.ahmed.kadar@outlook.com, abokor.ahmed.kadar.nour@gmail.com'
        mail_subject = 'HealthCheck'

        mail_message_body = """\
        <html>
        <head></head>
        <body>
            <p>Bonjour<br><br>
            Veuillez trouver ci-joint les stats journalier du TopUp<br><br>
            Vous trouverez plus de détails <a href="%s">ici</a><br><br>
            Cordialement,<br><br>
            XANA SYSTEM
            </p>
        </body>
        </html>
        """ % (link)

        message = MIMEMultipart('alternative')
        message['Subject'] = mail_subject
        message['From'] = mail_from
        message['To'] = mail_to
        pptx = MIMEApplication(open('./stats/daily/'+name,'rb').read())
        pptx.add_header('content-disposition', 'attachment',filename=name)
        message.attach(pptx)
        message.attach(MIMEText(mail_message_body, 'html'))
        # Sent Email
        server = smtplib.SMTP_SSL('smtp.gmail.com', 465)
        server.login(gmail_user, gmail_password)
        server.sendmail(mail_from,mail_to.split(","), message.as_string())
        server.close()

        #################################################################################

        message = "HealthCheck TopUp Done and send !"
        etat = "S"
        Logs.objects.create(date_time=datetime_now,process=PROCESS_NUMBER,message_description=message,etat=etat)